import sys
import json
from pathlib import Path

def analyze_pptx(pptx_path: Path):
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except Exception as e:
        return {"error": f"Missing dependency python-pptx: {e}"}

    if not pptx_path.exists():
        return {"error": f"File not found: {pptx_path}"}

    prs = Presentation(str(pptx_path))
    slides_summary = []
    for i, slide in enumerate(prs.slides, start=1):
        title = None
        bullets = []
        images = 0
        tables = 0
        charts = 0

        # Title
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip() or None

        # Collect text from all shapes
        for shp in slide.shapes:
            try:
                if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images += 1
                elif shp.shape_type == MSO_SHAPE_TYPE.TABLE:
                    tables += 1
                # Heuristic: charts are embedded OLE/graphic frames; count graphic frames
                elif shp.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shp, 'chart'):
                    charts += 1
                if hasattr(shp, 'has_text_frame') and shp.has_text_frame:
                    text = '\n'.join(p.text for p in shp.text_frame.paragraphs).strip()
                    if not text:
                        continue
                    # Avoid duplicating title
                    if title and text == title:
                        continue
                    # Split by lines to emulate bullets
                    for line in [l.strip() for l in text.splitlines() if l.strip()]:
                        bullets.append(line)
            except Exception:
                continue

        # Notes
        notes = None
        try:
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text
                notes = notes_text.strip() or None
        except Exception:
            pass

        slides_summary.append({
            "index": i,
            "title": title,
            "bullets": bullets[:20],  # cap for brevity
            "counts": {"images": images, "tables": tables, "charts": charts},
            "notes": notes,
        })

    return {
        "file": str(pptx_path),
        "slides": len(prs.slides),
        "slides_summary": slides_summary,
    }

def main():
    if len(sys.argv) < 2:
        print("Usage: python analyze_ppt.py <path-to-pptx>")
        sys.exit(2)
    pptx_path = Path(sys.argv[1])
    result = analyze_pptx(pptx_path)
    print(json.dumps(result, ensure_ascii=False, indent=2))

if __name__ == "__main__":
    main()
